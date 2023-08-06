import pptx.table
from pptx.util import Inches, Cm
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from typing import Union
from pptx.util import Emu
from pptx.oxml.shapes.graphfrm import CT_GraphicalObjectFrame
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.placeholder import PlaceholderGraphicFrame, PlaceholderPicture, ChartPlaceholder, TablePlaceholder
from pptx.oxml.shapes.picture import CT_Picture


class PlaceholderComponent:
    """Placeholder shape that can only accept a table."""

    @classmethod
    def insert_table(cls, self, rows, cols):
        """Return |PlaceholderGraphicFrame| object containing a `rows` by `cols` table.

        The position and width of the table are those of the placeholder and its height
        is proportional to the number of rows. A |PlaceholderGraphicFrame| object has
        all the properties and methods of a |GraphicFrame| shape except that the value
        of its :attr:`~._BaseSlidePlaceholder.shape_type` property is unconditionally
        `MSO_SHAPE_TYPE.PLACEHOLDER`. Note that the return value is not the new table
        but rather *contains* the new table. The table can be accessed using the
        :attr:`~.PlaceholderGraphicFrame.table` property of the returned
        |PlaceholderGraphicFrame| object.
        """
        graphicFrame = self._new_placeholder_table(rows, cols)
        self._replace_placeholder_with(graphicFrame)
        return PlaceholderGraphicFrame(graphicFrame, self._parent)

    @classmethod
    def _new_placeholder_table(cls, self, rows, cols):
        """
        Return a newly added `p:graphicFrame` element containing an empty
        table with *rows* rows and *cols* columns, positioned at the location
        of this placeholder and having its same width. The table's height is
        determined by the number of rows.
        """
        shape_id, name, height = self.shape_id, self.name, Emu(rows * 370840)
        return CT_GraphicalObjectFrame.new_table_graphicFrame(
            shape_id, name, rows, cols, self.left, self.top, self.width, height
        )

    @classmethod
    def insert_chart(cls, self, chart_type, chart_data):
        """
        Return a |PlaceholderGraphicFrame| object containing a new chart of
        *chart_type* depicting *chart_data* and having the same position and
        size as this placeholder. *chart_type* is one of the
        :ref:`XlChartType` enumeration values. *chart_data* is a |ChartData|
        object populated with the categories and series values for the chart.
        Note that the new |Chart| object is not returned directly. The chart
        object may be accessed using the
        :attr:`~.PlaceholderGraphicFrame.chart` property of the returned
        |PlaceholderGraphicFrame| object.
        """
        rId = self.part.add_chart_part(chart_type, chart_data)
        graphicFrame = self._new_chart_graphicFrame(
            rId, self.left, self.top, self.width, self.height
        )
        self._replace_placeholder_with(graphicFrame)
        return PlaceholderGraphicFrame(graphicFrame, self._parent)

    @classmethod
    def _new_chart_graphicFrame(cls, self, rId, x, y, cx, cy):
        """
        Return a newly created `p:graphicFrame` element having the specified
        position and size and containing the chart identified by *rId*.
        """
        id_, name = self.shape_id, self.name
        return CT_GraphicalObjectFrame.new_chart_graphicFrame(
            id_, name, rId, x, y, cx, cy
        )

    @classmethod
    def insert_picture(cls, self, image_file):
        """Return a |PlaceholderPicture| object depicting the image in `image_file`.

        `image_file` may be either a path (string) or a file-like object. The image is
        cropped to fill the entire space of the placeholder. A |PlaceholderPicture|
        object has all the properties and methods of a |Picture| shape except that the
        value of its :attr:`~._BaseSlidePlaceholder.shape_type` property is
        `MSO_SHAPE_TYPE.PLACEHOLDER` instead of `MSO_SHAPE_TYPE.PICTURE`.
        """
        pic = self._new_placeholder_pic(image_file)
        self._replace_placeholder_with(pic)
        return PlaceholderPicture(pic, self._parent)

    @classmethod
    def _new_placeholder_pic(cls, self, image_file):
        """
        Return a new `p:pic` element depicting the image in *image_file*,
        suitable for use as a placeholder. In particular this means not
        having an `a:xfrm` element, allowing its extents to be inherited from
        its layout placeholder.
        """
        rId, desc, image_size = self._get_or_add_image(image_file)
        shape_id, name = self.shape_id, self.name
        pic = CT_Picture.new_ph_pic(shape_id, name, desc, rId)
        pic.crop_to_fit(image_size, (self.width, self.height))
        return pic

    @classmethod
    def _get_or_add_image(cls, self, image_file):
        """
        Return an (rId, description, image_size) 3-tuple identifying the
        related image part containing *image_file* and describing the image.
        """
        image_part, rId = self.part.get_or_add_image_part(image_file)
        desc, image_size = image_part.desc, image_part._px_size
        return rId, desc, image_size


class PPTTable:
    def __init__(self, table: pptx.table):
        self.table = table

    @staticmethod
    def hex_to_rgb(color):
        return tuple(int(color[i:i + 2], 16) for i in (0, 2, 4))

    @staticmethod
    def _sub_element(parent, tagname, **kwargs):
        element = OxmlElement(tagname)
        element.attrib.update(kwargs)
        parent.append(element)
        return element

    def set_header_color(self, color):
        fill = RGBColor(*self.hex_to_rgb(color))
        for r in range(0, 1, 1):
            for cell in self.table.rows[r].cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill

    def set_odd_row_color(self, color):
        fill = RGBColor(*self.hex_to_rgb(color))
        for r in range(1, len(self.table.rows), 2):
            for cell in self.table.rows[r].cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill

    def set_even_row_color(self, color):
        fill = RGBColor(*self.hex_to_rgb(color))
        for r in range(2, len(self.table.rows), 2):
            for cell in self.table.rows[r].cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = fill

    def set_cell_border(self, border_color="000000", border_width='12700'):
        """ Hack function to enable the setting of border width and border color
            - bottom border only at present
            (c) Steve Canny
        """
        for r in range(0, len(self.table.rows), 1):
            for cell in self.table.rows[r].cells:
                tc = cell._tc
                tcPr = tc.get_or_add_tcPr()
                for lines in ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']:
                    lnR = self._sub_element(
                        tcPr, lines, w=border_width, cap='flat', cmpd='sng', algn='ctr')
                    solidFill = self._sub_element(lnR, 'a:solidFill')
                    srgbClr = self._sub_element(solidFill, 'a:srgbClr', val=border_color)
                cell.fill.solid()

