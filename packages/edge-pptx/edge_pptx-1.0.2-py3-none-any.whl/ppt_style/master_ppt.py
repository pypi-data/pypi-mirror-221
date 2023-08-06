import types
import pptx.table
import os
from typing import Union
from .component import PlaceholderComponent, PPTTable
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE
import requests
from io import BytesIO


class PPTGenerator:
    def __init__(self, template_ppt: Union[None, str, Presentation] = None):
        if type(template_ppt) == Presentation:
            self.presentation = template_ppt
        elif isinstance(template_ppt, str):
            self.presentation = Presentation(template_ppt)
        else:
            self.presentation = Presentation(self._default_pptx_path)

    @property
    def _default_pptx_path(self):
        """
        Return the path to the built-in default .pptx package.
        """
        _thisdir = os.path.split(__file__)[0]
        return os.path.join(_thisdir, "templates", "default.pptx")

    def get_placeholder(self, slide, placeholder_index):
        placeholder = slide.shapes[placeholder_index]
        left = placeholder.left
        top = placeholder.top
        width = placeholder.width
        height = placeholder.height
        return slide, placeholder, left, top, width, height

    def remove_placeholder(self, placeholder):
        sp = placeholder._sp
        sp.getparent().remove(sp)

    def replace_placeholder(self, slide, placeholder_index):
        slide, placeholder, left, top, width, height = self.get_placeholder(slide, placeholder_index)
        self.remove_placeholder(placeholder)
        return slide, left, top, width, height

    def insert_table(self, data, slide, placeholder_index):
        # slide, left, top, width, height = self.replace_placeholder(slide, placeholder_index)
        slide, placeholder, left, top, width, height = self.get_placeholder(slide, placeholder_index)
        stats = data
        # shape_table = slide.shapes.add_table(stats.shape[0] + 1, stats.shape[1], left, top, width, height).table
        placeholder._new_placeholder_table = types.MethodType(PlaceholderComponent._new_placeholder_table, placeholder)
        placeholder.insert_table = types.MethodType(PlaceholderComponent.insert_table, placeholder)
        shape_table = placeholder.insert_table(stats.shape[0] + 1, stats.shape[1]).table
        for i in range(stats.shape[1]):
            shape_table.cell(0, i).text = stats.columns[i]

        for i in range(stats.shape[0]):
            for j in range(stats.shape[1]):
                shape_table.cell(i + 1, j).text = str(stats.iloc[i, j])

        ppt_table = PPTTable(shape_table)
        ppt_table.set_even_row_color("FFFFFF")
        ppt_table.set_odd_row_color("DCE6F2")
        ppt_table.set_header_color("08409C")
        ppt_table.set_cell_border("BFBFBF")

    def insert_text(self, content_text, slide, placeholder_index):
        slide, placeholder, _, _, _, _ = self.get_placeholder(slide, placeholder_index)
        placeholder.text = content_text

    def insert_picture(self, file_path, slide, placeholder_index):
        slide, placeholder, _, _, _, _ = self.get_placeholder(slide, placeholder_index)
        placeholder._get_or_add_image = types.MethodType(PlaceholderComponent._get_or_add_image, placeholder)
        placeholder._new_placeholder_pic = types.MethodType(PlaceholderComponent._new_placeholder_pic, placeholder)
        placeholder.insert_picture = types.MethodType(PlaceholderComponent.insert_picture, placeholder)
        if "http" in file_path:
            response = requests.get(file_path)
            image_content = response.content
            image = BytesIO(image_content)
            placeholder.insert_picture(image)
        else:
            placeholder.insert_picture(file_path)

    def insert_chart(self, chart_type: XL_CHART_TYPE, chart_data, slide, placeholder_index):
        slide, placeholder, _, _, _, _ = self.get_placeholder(slide, placeholder_index)
        placeholder._new_chart_graphicFrame = types.MethodType(PlaceholderComponent._new_chart_graphicFrame,
                                                               placeholder)
        placeholder.insert_chart = types.MethodType(PlaceholderComponent.insert_chart, placeholder)
        placeholder.insert_chart(chart_type, chart_data)

    def save_ppt(self, output_file):
        self.presentation.save(output_file)
